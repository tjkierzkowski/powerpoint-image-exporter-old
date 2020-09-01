import pytest
from pathlib import Path
from pptx import Presentation

from pptx_export.pptx_export import PowerPointImageExporter


@pytest.fixture
def fake_path(tmp_path):
    fake_file_path = tmp_path / "presentations"
    return fake_file_path


@pytest.fixture
def custom_path(tmp_path: Path) -> Path:
    custom_path = tmp_path / "presentation_images"
    return custom_path


@pytest.fixture
def default_path(tmp_path: Path) -> Path:
    custom_path = tmp_path / "lecture_images"
    return custom_path


@pytest.fixture
def fake_file():
    return 'a_valid_file.pptx'


@pytest.fixture
def minimal_pres(tmp_path):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"
    minimal_file = tmp_path / 'minimal.pptx'
    prs.save(minimal_file)
    return minimal_file


@pytest.fixture
def valid_presentation_name(custom_path, fake_file):
    ppt_stub_file = custom_path / fake_file
    pres = PowerPointImageExporter(ppt_stub_file)
    return pres


# Integration specific fixtures below

@pytest.fixture
def default_path(tmp_path: Path) -> Path:
    custom_path = tmp_path / "lecture_images"
    return custom_path


@pytest.fixture(scope="session")
def actual_presentation_path():
    """Load an actual .pptx file with images to test against"""
    project_name = 'powerpoint_image_exporter'
    presentation_under_test = 'stub_tester.pptx'
    project_dir = [project for project in Path.home().rglob(project_name) if project.is_dir()]
    if not project_dir:
        raise ValueError(f"Could not find project directory '{project_name}' from your user's home directory")
    project_root = project_dir[0]
    presentations = [pres for pres in project_root.rglob(presentation_under_test) if pres.is_file()]
    if not presentations:
        raise ValueError(f"Could not find the actual presentation '{presentation_under_test}' within the project "
                         f"directory {project_root.resolve()}")
    return str(presentations[0].resolve())