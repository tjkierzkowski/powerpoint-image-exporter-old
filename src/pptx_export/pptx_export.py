"""
export images from a pptx file
"""
import logging
from pathlib import Path

# from pip._internal.utils.deprecation import deprecated
from pptx import Presentation
from pptx.shapes.autoshape import Shape
from typing import Tuple
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.placeholder import PlaceholderPicture

DEFAULT_DIR = "lecture_images"


class PowerPointImageExporter:
    """All-in-one class to copy out images from a single PowerPoint file into a directory"""

    def __init__(self, pptx_file):
        if pptx_file is None:
            raise ValueError("Please provide a valid PowerPoint file")
        pptx_file_path = Path(pptx_file)
        if pptx_file_path.exists() and not pptx_file_path.is_file():
            raise ValueError("The provided file path is something other than a file")
        if pptx_file_path.suffixes[-1] != ".pptx":
            raise ValueError(
                f"Cannot export images from something other than a .pptx file. The file extension entered "
                f"was: {''.join(pptx_file_path.suffixes)}"
            )
        self.pptx_file = pptx_file_path
        self.safe_presentation_name = self.__snakecase_the_ppt_name(self.pptx_file.stem)
        self.image_directory_path = None
        self.default_image_path = str(Path("").resolve() / DEFAULT_DIR)

    def __snakecase_the_ppt_name(self, filename):
        return filename.replace(" ", "_")

    def create_directory_for_images(self, new_directory=None):
        """Create a directory if none exists and skips creation if the directory provided already exists. If the
        new_directory is set to None a ValueError is thrown."""
        if new_directory is None:
            raise ValueError(
                f"This is not a valid directory path for the images to be placed in: {new_directory}"
            )
        image_directory_path = Path(new_directory)
        if image_directory_path.exists() and image_directory_path.is_dir():
            logging.info(
                f"directory '{image_directory_path.resolve()}' already exists... skipping creation"
            )
        else:
            image_directory_path.mkdir()
            logging.info(f"directory created at {image_directory_path.resolve()}")
        self.image_directory_path = image_directory_path

    def copy_images_to_directory(self):
        """Copies all of the images from the powerpoint file into an empty directory.  """
        if self.image_directory_path is None:
            self.create_directory_for_images(self.default_image_path)
        if self.image_directory_path.exists() and self.image_directory_path.is_dir():
            image_directory_contained_files = [
                images for images in self.image_directory_path.iterdir()
            ]
            if len(image_directory_contained_files) != 0:
                raise ValueError(
                    f"Will not overwrite the existing image directory at: {self.image_directory_path.resolve()}"
                )

        presentation = Presentation(self.pptx_file)

        for picture_info in self.iter_by_shape(presentation, MSO_SHAPE_TYPE.PICTURE):
            slide_number, image_number, picture = picture_info
            ext = picture.image.ext
            image_bytes = picture.image.blob
            image_filename = self.meaningful_filename(
                self.safe_presentation_name, slide_number, image_number, ext
            )
            full_image_filename = self.image_directory_path / image_filename
            with open(full_image_filename, "wb") as img_out:
                img_out.write(image_bytes)

    def iter_by_shape(
        self, presentation: Presentation, shape_type: MSO_SHAPE_TYPE
    ) -> Tuple[int, int, Shape]:
        """Generator function for returning a Picture Shape or PlaceholderPicture along with the slide number and shape
        number in each slide in the powerpoint"""
        if presentation is None:
            raise ValueError(
                f"Presentation {presentation} cannot be traversed to extract PowerPoint elements"
            )
        if shape_type is None:
            raise ValueError(
                f"Please provide a valid Shape enum to traverse your presentation, see pptx docs for "
                f"examples"
            )
        for slide_number, slide in enumerate(presentation.slides):
            image_number = 0
            for shape in slide.shapes:
                if (shape.is_placeholder and isinstance(shape, PlaceholderPicture)) or (
                    shape.shape_type == shape_type
                ):
                    image_number += 1
                    yield slide_number + 1, image_number, shape

    def meaningful_filename(self, pres_name, slide_number, image_number, ext) -> str:
        """generate an image file name in the form of:
        <ppt_file_name>_<slide_number>_<image_number>
        where:
        ppt_file_name: is the name of the .pptx lecture without spaces
        slide_number: is the associated slide number of the file
        image_number: is the image number relative to the slide
        """
        return "{}_slide_{}_image_{}.{}".format(
            pres_name, slide_number, image_number, ext
        )
