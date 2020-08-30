"""
export images from a pptx file
"""
from pathlib import Path
from pptx import Presentation
from pptx.shapes.autoshape import Shape
from typing import Tuple, Any
from pptx.enum.shapes import MSO_SHAPE_TYPE


class PowerPointImageExporter:

    def __init__(self, pptx_file):
        if pptx_file is None:
            raise ValueError("Please provide a valid PowerPoint file")
        pptx_file_path = Path(pptx_file)
        if pptx_file_path.exists() and not pptx_file_path.is_file():
            raise ValueError("The provided file path is something other than a file")
        if pptx_file_path.suffixes[-1] != '.pptx':
            raise ValueError(f"Cannot export images from something other than a .pptx file. The file extension entered "
                             f"was: {''.join(pptx_file_path.suffixes)}")
        self.pptx_file = pptx_file_path
        self.safe_presentation_name = self._snakecase_the_ppt_name(self.pptx_file.stem)
        self.image_directory_path = None
        default_dir = 'lecture_images'
        self._default_image_path = str( Path('.') / default_dir)

    def _snakecase_the_ppt_name(self, filename):
        return filename.replace(" ", "_")

    def __find_desktop_directory_path(self):
        desktop_paths = [desktop for desktop in Path.home().glob('Desktop') if desktop.is_dir()]
        return Path(desktop_paths[0])

    def create_directory_for_images(self, new_directory=None):
        """Create a directory if none exists. Also sets the image_directory_path to the new path. If the
        new_directory is set to None a ValueError is thrown. """
        if new_directory is None:
            raise ValueError(f"This is not a valid directory path for the images to be placed in: {new_directory}")
        image_directory_path = Path(new_directory)
        if image_directory_path.exists() and image_directory_path.is_dir():
            print(f"directory '{image_directory_path.resolve()}' already exists... skipping creation\n")
        else:
            image_directory_path.mkdir()
            print(f"directory created at {image_directory_path.resolve()}\n")
        self.image_directory_path = image_directory_path

    def copy_images_to_directory(self):
        print(f'opening {self.pptx_file.name}')
        presentation = Presentation(self.pptx_file)
        if self.image_directory_path is None:
            self.create_directory_for_images(self._default_image_path)
        elif self.image_directory_path.exists() and self.image_directory_path.is_dir():
            contained_files = [images for images in self.image_directory_path.iterdir()]
            if len(contained_files) != 0:
                raise ValueError(f"Will not overwrite the existing image directory at: {self.image_directory_path.resolve()}")

        for picture_info in self.iter_by_shape(presentation, MSO_SHAPE_TYPE.PICTURE):
            slide_number, image_number, picture = picture_info
            ext = picture.image.ext
            image_bytes = picture.image.blob
            image_filename = self.meaningful_filename(self.safe_presentation_name, slide_number, image_number, ext)
            full_image_filename = self.image_directory_path / image_filename
            with open(full_image_filename, 'wb') as img_out:
                img_out.write(image_bytes)

    def iter_by_shape(self, presentation: Presentation, shape_type: MSO_SHAPE_TYPE) -> Tuple[int, int, Shape]:
        if presentation is None:
            raise ValueError(f"Presentation {presentation} cannot be traversed to extract PowerPoint elements")
        if shape_type is None:
            raise ValueError(f"Please provide a valid Shape enum to traverse your presentation, see pptx docs for "
                             f"examples")
        for slide_number, slide in enumerate(presentation.slides):
            for shape_number, shape in enumerate(slide.shapes):
                if shape.shape_type == shape_type:
                    yield slide_number, shape_number, shape

    def meaningful_filename(self, pres_name, slide_number, image_number, ext) -> str:
        """generate an image file name in the form of:
            <ppt_file_name>_<slide_number>_<image_number>
            where:
            ppt_file_name: is the name of the file
            slide_number: is the slide number of the file
            image_number: the image number relative to the slide
        """
        return "{}_slide_{}_image_{}.{}".format(pres_name, slide_number, image_number, ext)
