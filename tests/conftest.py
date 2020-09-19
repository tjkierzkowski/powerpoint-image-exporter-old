from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
import pytest

from pptx_export.pptx_export import DEFAULT_DIR, PowerPointImageExporter


@pytest.fixture
def fake_path(tmp_path):
    fake_file_path = tmp_path / "presentations"
    return fake_file_path


@pytest.fixture
def custom_path(tmp_path: Path) -> Path:
    custom_path = tmp_path / "presentation_images"
    return custom_path


@pytest.fixture(scope="session")
def default_path_session(tmp_path_factory) -> Path:
    base = tmp_path_factory.mktemp()
    custom_path = base / DEFAULT_DIR
    return custom_path


@pytest.fixture
def fake_file():
    return "a_valid_file.pptx"


@pytest.fixture
def valid_presentation_name(custom_path, fake_file):
    ppt_stub_file = custom_path / fake_file
    pres = PowerPointImageExporter(ppt_stub_file)
    return pres


# Integration specific fixtures below


@pytest.fixture
def default_path(tmp_path: Path) -> Path:
    current_default = tmp_path / DEFAULT_DIR
    return current_default


@pytest.fixture(scope="session")
def data_directory(pytestconfig):
    """Provides easy access to the test data"""
    project_root = Path(pytestconfig.rootdir)
    return project_root / "data"


@pytest.fixture(scope="session")
def minimal_pres(data_directory):
    """Generates a PowerPoint presentation based on the random images placed in the
    data directory.  Creates placeholder images as well as actual picture shapes.
    GroupShapes, or images grouped together in a slide are not generated yet.
    """

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Fake Presentation"
    subtitle.text = "Another presentation brought to you by python-pptx!"
    # Add the images from the data directory into a presentation
    legal_image_types = {".jpg", ".jpeg", ".gif", ".tiff", ".png"}
    image_types = 2
    for image_num, image_file in enumerate(data_directory.iterdir()):
        if (
            image_file
            and image_file.is_file()
            and image_file.suffix in legal_image_types
        ):
            image_file_path = str(image_file.resolve())
            if image_num % image_types == 0:
                # for adding a picture per slide as a shape
                blank_slide_layout = prs.slide_layouts[6]
                picture_slide = prs.slides.add_slide(blank_slide_layout)
                picture_slide.shapes.add_picture(image_file_path, Inches(1), Inches(2))
            else:
                # for adding a picture per slide as a placeholder
                slide = prs.slides.add_slide(prs.slide_layouts[8])
                placeholder = slide.placeholders[1]
                placeholder.insert_picture(image_file_path)

    minimal_file = data_directory / "minimal.pptx"
    prs.save(str(minimal_file.resolve()))
    yield minimal_file
    minimal_file.unlink(missing_ok=True)
