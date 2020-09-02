import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.placeholder import PlaceholderPicture
from pptx_export.pptx_export import PowerPointImageExporter


def test_powerpoint_images_are_the_same_as_the_number_in_the_file(
    tmp_path, minimal_pres
):
    custom_path = tmp_path / "presentation_images"
    pres_file = minimal_pres
    pptx_export = PowerPointImageExporter(pres_file)
    pres = Presentation(pres_file)
    result_pictures = [
        picture_info
        for picture_info in pptx_export.iter_by_shape(pres, MSO_SHAPE_TYPE.PICTURE)
    ]
    result_amount = len(result_pictures)
    pptx_export.create_directory_for_images(custom_path)
    expected_images = [
        shape
        for slide in pres.slides
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        or shape.is_placeholder
        and isinstance(shape, PlaceholderPicture)
    ]
    assert len(expected_images) == result_amount


def test_if_all_images_from_an_actual_presentation_are_extracted_to_directory(
    default_path, minimal_pres
):
    real_presentation = minimal_pres
    pptx_export = PowerPointImageExporter(real_presentation)
    pptx_export.create_directory_for_images(default_path)
    pptx_export.copy_images_to_directory()
    result_pictures = [files for files in default_path.iterdir()]
    pres = Presentation(real_presentation)
    expected_images = [
        shape
        for slide in pres.slides
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        or shape.is_placeholder
        and isinstance(shape, PlaceholderPicture)
    ]
    assert len(expected_images) == len(result_pictures) and len(expected_images) > 0
